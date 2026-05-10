#!/usr/bin/env python3
"""
Generate 3-slide before/after comparison PowerPoint.
Same design system as presentations/generate_research_pptx.py.

Run:  python generate_comparison_pptx.py
Out:  pipeline_v1_vs_v2.pptx
"""

import io
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette (identical to existing presentations) ─────────────────────────────
BLUE     = "#0072B2"
CHARCOAL = "#333333"
DARKGRAY = "#666666"
GRAY     = "#AAAAAA"
LIGHTGRAY= "#DDDDDD"
LIGHT    = "#F5F5F5"
BG       = "#FFFFFF"
RED      = "#D55E00"   # for "bad" / v1 problems

def rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

plt.rcParams.update({
    "font.family": "sans-serif",
    "font.sans-serif": ["Inter", "Helvetica Neue", "Arial"],
    "axes.spines.top": False, "axes.spines.right": False,
    "axes.labelcolor": CHARCOAL, "xtick.color": DARKGRAY,
    "ytick.color": DARKGRAY, "text.color": CHARCOAL,
    "figure.facecolor": BG, "axes.facecolor": BG,
})

FONT = "Helvetica Neue"
W = Inches(13.33)
H = Inches(7.5)
ML = Inches(0.9)
MT = Inches(0.75)
CW = Inches(11.53)
TOTAL_SLIDES = 3


# ─────────────────────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────────────────────

def fig_to_stream(fig, dpi=180):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight",
                facecolor=BG, edgecolor="none")
    buf.seek(0)
    plt.close(fig)
    return buf

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_rect(slide, left, top, width, height, fill_hex):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_hex)
    return shape

def add_textbox(slide, left, top, width, height, text,
                font_size=18, bold=False, italic=False,
                color=CHARCOAL, align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    return txBox

def add_multiline(slide, left, top, width, height, lines, font_size=13, color=CHARCOAL):
    """Add a textbox with multiple paragraphs (list of (text, bold, color) tuples)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, bld, clr) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(font_size)
        run.font.bold = bld
        run.font.color.rgb = rgb(clr)
    return txBox

def add_image(slide, stream, left, top, width=None, height=None):
    stream.seek(0)
    return slide.shapes.add_picture(stream, left, top, width=width, height=height)

def add_slide_number(slide, n):
    add_textbox(slide, W - Inches(1.2), H - Inches(0.42), Inches(1.0), Inches(0.3),
                f"{n} / {TOTAL_SLIDES}", font_size=11, color=GRAY, align=PP_ALIGN.RIGHT)

def add_title(slide, text, font_size=24):
    add_textbox(slide, ML, MT, CW, Inches(0.9),
                text, font_size=font_size, bold=True, color=CHARCOAL)
    add_rect(slide, ML, Inches(1.72), Inches(0.6), Inches(0.045), BLUE)


# ─────────────────────────────────────────────────────────────────────────────
# Figures
# ─────────────────────────────────────────────────────────────────────────────

def fig_params_comparison():
    """Bar chart: params per positive sample, v1 vs v2."""
    fig, ax = plt.subplots(figsize=(5.5, 3.8))
    labels = ["v1\nSwin UNETR", "v2\nResNet-10"]
    params = [62.0, 14.5]
    positives = [400, 1040]
    ratio = [p * 1e6 / n for p, n in zip(params, positives)]

    colors = [LIGHTGRAY, BLUE]
    bars = ax.bar(labels, ratio, color=colors, width=0.5, edgecolor="none")

    for bar, r, p, n in zip(bars, ratio, params, positives):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 2000,
                f"{r:,.0f}", ha="center", fontsize=14, fontweight="bold",
                color=RED if r > 100000 else BLUE)
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() * 0.5,
                f"{p:.0f}M params\n{n} positives", ha="center", fontsize=10,
                color=CHARCOAL if r > 100000 else BG)

    ax.set_ylabel("Parameters per positive sample", fontsize=12)
    ax.spines["bottom"].set_color(LIGHTGRAY)
    ax.spines["left"].set_color(LIGHTGRAY)
    ax.yaxis.grid(True, color="#F0F0F0", zorder=0)
    ax.set_axisbelow(True)

    # Red danger zone
    ax.axhline(50000, color=RED, lw=1.2, ls="--", alpha=0.6)
    ax.text(1.35, 52000, "overfitting zone", fontsize=9, color=RED, style="italic")
    fig.tight_layout()
    return fig_to_stream(fig)


def fig_training_strategy():
    """Visual diagram of 2-stage training."""
    fig, ax = plt.subplots(figsize=(10, 2.8))
    ax.set_xlim(0, 40)
    ax.set_ylim(0, 3.2)
    ax.axis("off")

    # Stage 1: CE warmup
    rect1 = mpatches.FancyBboxPatch(
        (0.5, 0.8), 14, 1.8, boxstyle="round,pad=0.12",
        facecolor=LIGHT, edgecolor=LIGHTGRAY, linewidth=1.5, zorder=3)
    ax.add_patch(rect1)
    ax.text(7.5, 2.1, "STAGE 1  —  Epochs 1-15", ha="center",
            fontsize=11, fontweight="bold", color=CHARCOAL, zorder=4)
    ax.text(7.5, 1.45, "BCE + Focal Loss (warmup)\nAdamW, layer-wise LR decay\nFreeze backbone ep 1-5",
            ha="center", fontsize=9.5, color=DARKGRAY, zorder=4, linespacing=1.5)

    # Arrow
    ax.annotate("", xy=(16, 1.7), xytext=(14.8, 1.7),
                arrowprops=dict(arrowstyle="-|>", color=BLUE, lw=2.5), zorder=5)

    # Stage 2: APLoss
    rect2 = mpatches.FancyBboxPatch(
        (16.5, 0.8), 14, 1.8, boxstyle="round,pad=0.12",
        facecolor=BLUE, edgecolor=BLUE, linewidth=1.5, zorder=3)
    ax.add_patch(rect2)
    ax.text(23.5, 2.1, "STAGE 2  —  Epochs 16-40", ha="center",
            fontsize=11, fontweight="bold", color=BG, zorder=4)
    ax.text(23.5, 1.45, "APLoss + SOAP optimizer\nDualSampler (1 pos/batch)\nCosine LR decay",
            ha="center", fontsize=9.5, color="#CCE5F5", zorder=4, linespacing=1.5)

    # Labels below
    ax.text(7.5, 0.3, "Gives the model a reasonable score distribution",
            ha="center", fontsize=9, color=GRAY, style="italic")
    ax.text(23.5, 0.3, "Directly optimizes AUPREC (ranking metric)",
            ha="center", fontsize=9, color=GRAY, style="italic")

    # v1 label
    ax.text(33, 2.6, "v1: APLoss from epoch 1\n(no warmup)", fontsize=9,
            color=RED, style="italic", ha="center")
    ax.annotate("", xy=(30.5, 2.1), xytext=(32, 2.4),
                arrowprops=dict(arrowstyle="->", color=RED, lw=1.0), zorder=2)

    fig.tight_layout()
    return fig_to_stream(fig)


def fig_changes_impact():
    """Horizontal bar chart of expected impact per change."""
    fig, ax = plt.subplots(figsize=(8, 3.5))

    changes = [
        "Model 14.5M (was 62M)",
        "2.6x more positives (1040 vs 400)",
        "CE warmup before APLoss",
        "LR scheduler fix",
        "T2/T1 ratio 3rd channel",
        "Layer-wise LR decay",
        "EMA weights",
        "Gradient accumulation (bs 16)",
    ]
    impact = [5, 4, 4, 3, 3, 2, 2, 2]

    colors_list = [BLUE if v >= 4 else LIGHTGRAY for v in impact]
    y = np.arange(len(changes))
    bars = ax.barh(y, impact, color=colors_list, height=0.55, edgecolor="none")

    ax.set_yticks(y)
    ax.set_yticklabels(changes, fontsize=11)
    ax.set_xlabel("Expected impact (1-5)", fontsize=11)
    ax.set_xlim(0, 6)
    ax.invert_yaxis()
    ax.spines["bottom"].set_color(LIGHTGRAY)
    ax.spines["left"].set_color(LIGHTGRAY)
    ax.xaxis.grid(True, color="#F0F0F0", zorder=0)
    ax.set_axisbelow(True)

    for bar, v in zip(bars, impact):
        label = ["", "Low", "Moderate", "High", "Critical", "Critical"][v]
        ax.text(bar.get_width() + 0.1, bar.get_y() + bar.get_height()/2,
                label, va="center", fontsize=10,
                color=BLUE if v >= 4 else DARKGRAY, fontweight="bold" if v >= 4 else "normal")
    fig.tight_layout()
    return fig_to_stream(fig)


# ─────────────────────────────────────────────────────────────────────────────
# Slides
# ─────────────────────────────────────────────────────────────────────────────

def slide_1_diagnostic(prs, img_params):
    """Why v1 is stuck at AUPREC 0.16."""
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, W, Inches(0.055), BLUE)
    add_title(slide, "Pipeline v1 is stuck at AUPREC 0.16 — the model memorizes, it doesn't learn")

    # Left: 4 root causes
    causes = [
        ("62M parameters for ~400 positives",
         "155,000 params per positive sample. The model memorizes training subjects "
         "instead of learning WMH features. Rule of thumb: <1,000 params/positive."),
        ("APLoss without CE warmup",
         "APLoss needs a reasonable starting score distribution. Without warmup, "
         "the surrogate loss has nothing to rank — gradients are near-random."),
        ("LR scheduler bug",
         "Cosine annealing stepped per batch but T_max set in epochs. "
         "LR collapsed to zero after a few batches instead of decaying over epochs."),
        ("CT abdominal SSL pretraining",
         "MONAI's Swin UNETR weights come from CT abdominal scans. "
         "The domain gap to pediatric brain MRI is enormous — effectively random init."),
    ]

    for i, (title, body) in enumerate(causes):
        y = Inches(1.92) + Inches(i * 1.15)
        # Red left border
        add_rect(slide, ML, y, Inches(0.05), Inches(0.95), RED)
        add_textbox(slide, ML + Inches(0.18), y + Inches(0.02), Inches(6.0), Inches(0.3),
                    f"{i+1}. {title}", font_size=13, bold=True, color=RED)
        add_textbox(slide, ML + Inches(0.18), y + Inches(0.38), Inches(6.0), Inches(0.55),
                    body, font_size=11, color=DARKGRAY)

    # Right: params chart
    add_image(slide, img_params, Inches(7.6), Inches(1.95), width=Inches(5.3))

    # Bottom: result
    add_rect(slide, ML, Inches(6.6), CW, Inches(0.5), LIGHT)
    add_textbox(slide, ML + Inches(0.15), Inches(6.65), CW, Inches(0.4),
                "Result: val AUPREC = 0.16 — the model produces near-random rankings on held-out subjects",
                font_size=13, bold=True, color=RED)

    add_slide_number(slide, 1)


def slide_2_comparison(prs, img_strategy):
    """Before vs After table + training strategy diagram."""
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, W, Inches(0.055), BLUE)
    add_title(slide, "v2 addresses every diagnosed failure mode with evidence-based changes")

    # Table header
    th = Inches(1.88)
    tw = Inches(11.53)
    col_w = [Inches(2.8), Inches(4.0), Inches(4.0), Inches(0.73)]
    col_x = [ML]
    for w in col_w[:-1]:
        col_x.append(col_x[-1] + w)

    add_rect(slide, ML, th, tw, Inches(0.38), CHARCOAL)
    headers = ["", "v1 (before)", "v2 (after)", ""]
    h_colors = [BG, BG, BG, BG]
    for header, x, w, c in zip(headers, col_x, col_w, h_colors):
        add_textbox(slide, x + Inches(0.08), th + Inches(0.05), w, Inches(0.3),
                    header, font_size=11, bold=True, color=c)

    rows = [
        ("Model",        "Swin UNETR  62M params",    "ResNet-10  14.5M params",    "11x"),
        ("Positives",    "~400 (1 CSV, score >= 3)",   "1,040 (2 CSV, WMA match)",   "2.6x"),
        ("Input",        "2ch (T1 + T2)",              "3ch (T1 + T2 + T2/T1 ratio)",""),
        ("Loss warmup",  "None — APLoss from epoch 1", "15 ep BCE+Focal, then APLoss",""),
        ("LR schedule",  "Cosine/batch (bugged)",      "Cosine/epoch (fixed)",       ""),
        ("LR unfreezing","Uniform LR/10",              "Layer-wise decay (0.75/layer)",""),
        ("EMA",          "No",                         "Yes (decay 0.9995)",         ""),
        ("Grad accum",   "No (bs=4)",                  "Yes (effective bs=16)",      "4x"),
    ]

    for i, (metric, v1, v2, gain) in enumerate(rows):
        y = th + Inches(0.38 + i * 0.42)
        bg = LIGHT if i % 2 == 0 else BG
        add_rect(slide, ML, y, tw, Inches(0.42), bg)
        add_textbox(slide, col_x[0] + Inches(0.08), y + Inches(0.06), col_w[0], Inches(0.32),
                    metric, font_size=12, bold=True, color=CHARCOAL)
        add_textbox(slide, col_x[1] + Inches(0.08), y + Inches(0.06), col_w[1], Inches(0.32),
                    v1, font_size=11, color=GRAY)
        add_textbox(slide, col_x[2] + Inches(0.08), y + Inches(0.06), col_w[2], Inches(0.32),
                    v2, font_size=11, bold=True, color=BLUE)
        if gain:
            add_textbox(slide, col_x[3] + Inches(0.08), y + Inches(0.06), col_w[3], Inches(0.32),
                        gain, font_size=11, bold=True, color=BLUE)

    # Training strategy diagram below table
    add_image(slide, img_strategy, ML, Inches(5.38), width=Inches(11.0))

    add_slide_number(slide, 2)


def slide_3_expected(prs, img_impact):
    """Expected impact + realistic targets."""
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, W, Inches(0.055), BLUE)
    add_title(slide, "Realistic target: AUPREC 0.35–0.55 — every change is literature-backed")

    # Left: impact chart
    add_image(slide, img_impact, ML, Inches(1.88), width=Inches(7.0))

    # Right: target box
    rl = Inches(8.3)
    rw = Inches(4.5)

    # Target box
    add_rect(slide, rl, Inches(1.88), rw, Inches(1.6), LIGHT)
    add_rect(slide, rl, Inches(1.88), rw, Inches(0.04), BLUE)
    add_textbox(slide, rl + Inches(0.15), Inches(2.0), rw - Inches(0.3), Inches(0.35),
                "AUPREC Targets", font_size=13, bold=True, color=BLUE)
    targets = [
        ("0.16", "v1 (current)", GRAY),
        ("0.35–0.55", "v2 realistic goal", BLUE),
        ("> 0.55", "publishable surprise", DARKGRAY),
        ("> 0.65", "suspect data leakage", RED),
    ]
    for i, (val, desc, clr) in enumerate(targets):
        y = Inches(2.45) + Inches(i * 0.25)
        add_textbox(slide, rl + Inches(0.2), y, Inches(1.2), Inches(0.25),
                    val, font_size=12, bold=True, color=clr)
        add_textbox(slide, rl + Inches(1.5), y, Inches(2.8), Inches(0.25),
                    desc, font_size=11, color=clr)

    # Early warning box
    add_rect(slide, rl, Inches(3.7), rw, Inches(1.75), LIGHT)
    add_rect(slide, rl, Inches(3.7), Inches(0.05), Inches(1.75), RED)
    add_textbox(slide, rl + Inches(0.18), Inches(3.78), rw - Inches(0.3), Inches(0.3),
                "Early Warning Signals", font_size=12, bold=True, color=RED)
    warnings = [
        "Stage 1 ep 15 AUPREC < 0.18 -> stop, debug",
        "Train AP > 0.40, Val < 0.20 -> overfitting",
        "Per-site variance > 0.15 -> site shortcut",
    ]
    for i, w in enumerate(warnings):
        add_textbox(slide, rl + Inches(0.18), Inches(4.15) + Inches(i * 0.38),
                    rw - Inches(0.3), Inches(0.35),
                    w, font_size=10, color=DARKGRAY)

    # Bottom: key message
    add_rect(slide, ML, Inches(5.9), CW, Inches(1.15), BLUE)
    add_textbox(slide, ML + Inches(0.2), Inches(5.98), CW - Inches(0.4), Inches(0.35),
                "KEY INSIGHT", font_size=11, bold=True, color="#99CCEE")
    add_textbox(slide, ML + Inches(0.2), Inches(6.35), CW - Inches(0.4), Inches(0.6),
                "The single biggest lever is the model size: 14.5M params on 1,040 positives "
                "(14k params/pos) vs 62M on 400 (155k params/pos). Combined with CE warmup "
                "and the LR fix, this addresses the three root causes of v1's failure.",
                font_size=14, bold=True, color=BG)

    add_slide_number(slide, 3)


# ─────────────────────────────────────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────────────────────────────────────

def main():
    print("Generating figures...")
    img_params   = fig_params_comparison()
    img_strategy = fig_training_strategy()
    img_impact   = fig_changes_impact()

    print("Building PowerPoint...")
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide_1_diagnostic(prs, img_params)
    slide_2_comparison(prs, img_strategy)
    slide_3_expected(prs, img_impact)

    out_path = Path(__file__).parent / "pipeline_v1_vs_v2.pptx"
    prs.save(str(out_path))
    print(f"Saved -> {out_path}")


if __name__ == "__main__":
    main()
